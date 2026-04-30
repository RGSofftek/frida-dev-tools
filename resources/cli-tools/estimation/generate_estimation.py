"""
Generador de presentaciones de estimacion RPA.

Uso:
    python generate_estimation.py --config estimacion.json [--out propuesta.pptx]

Las rutas de salida relativas se escriben en la carpeta output/ del proyecto.
"""

import argparse
import json
import os
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

TEMPLATE_DIR = Path(__file__).parent / "template"
OUTPUT_DIR = Path(__file__).resolve().parent / "output"


def resolve_output_path(out_arg: str) -> Path:
    """Rutas relativas bajo output/; rutas absolutas sin cambiar."""
    p = Path(out_arg)
    if p.is_absolute():
        return p
    return OUTPUT_DIR / p


def resolve_template_path() -> Path:
    """Busca template.pptx o RPA_Template.pptx en template/."""
    for name in ("template.pptx", "RPA_Template.pptx"):
        p = TEMPLATE_DIR / name
        if p.is_file():
            return p
    return TEMPLATE_DIR / "RPA_Template.pptx"

VALID_COMPLEXITY_LEVELS = ("Simple", "Medium", "High", "Very High")

TASK_KEYS = (
    "access_permits",
    "process_analysis",
    "process_doc",
    "env_config",
    "access_validation",
    "automation_dev",
    "test_data",
    "unit_test",
    "uat",
    "manual",
    "rollout",
    "monitoring",
)

# ── Colores ──────────────────────────────────────────────────────────────────

NAVY = RGBColor(0x12, 0x30, 0x4C)
BLUE = RGBColor(0x36, 0x5F, 0xA8)
CYAN = RGBColor(0x4C, 0xAB, 0xC0)
SECTION_FILL = RGBColor(0xC3, 0xCA, 0xD4)
ALT1 = RGBColor(0xF1, 0xF1, 0xF1)
ALT2 = RGBColor(0xE9, 0xE9, 0xE9)
DARK_TEXT = RGBColor(0x25, 0x2F, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_run_font(run, name="Nunito", size=None, bold=None, color=None):
    run.font.name = name
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color


def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def style_cell(cell, text="", fill=None, font_size=12, bold=False,
               color=DARK_TEXT, align=PP_ALIGN.LEFT):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill or WHITE
    cell.margin_left = Emu(110000)
    cell.margin_right = Emu(65000)
    cell.margin_top = Emu(25000)
    cell.margin_bottom = Emu(25000)
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            set_run_font(r, size=font_size, bold=bold, color=color)


def set_text_shape(shape, lines, title_size=24, body_size=16):
    """Rellena un shape con titulo (primera linea) y viñetas."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = lines[0]
    set_run_font(r, size=title_size, bold=True, color=NAVY)
    p.alignment = PP_ALIGN.LEFT
    for line in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        set_run_font(r, size=body_size, bold=False, color=NAVY)
        p.alignment = PP_ALIGN.LEFT
    tf.word_wrap = True
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)


def currency(v):
    return f"${v:,}"


# ── Validacion ───────────────────────────────────────────────────────────────

def _require(data, key_path):
    """Navega key_path (ej. 'cover.year') y falla si no existe."""
    keys = key_path.split(".")
    current = data
    for k in keys:
        if not isinstance(current, dict) or k not in current:
            print(f'Error: falta la clave "{key_path}" en el config.', file=sys.stderr)
            sys.exit(1)
        current = current[k]
    return current


def validate_config(data):
    _require(data, "cover.year")
    _require(data, "cover.automation_name")

    _require(data, "understanding.objectives")
    _require(data, "understanding.systems_tools")
    flow_img = _require(data, "understanding.flowchart_image")
    if not os.path.isfile(flow_img):
        print(f'Error: no se encuentra la imagen "{flow_img}".', file=sys.stderr)
        sys.exit(1)

    level = _require(data, "complexity.selected_level")
    if level not in VALID_COMPLEXITY_LEVELS:
        print(
            f'Error: complexity.selected_level "{level}" no es valido. '
            f"Valores aceptados: {', '.join(VALID_COMPLEXITY_LEVELS)}.",
            file=sys.stderr,
        )
        sys.exit(1)

    weeks = _require(data, "development_plan.weeks")
    tasks = _require(data, "development_plan.tasks")
    for key in TASK_KEYS:
        if key not in tasks:
            print(
                f'Error: falta la tarea "{key}" en development_plan.tasks. '
                f"Usa null si no aplica.",
                file=sys.stderr,
            )
            sys.exit(1)
        span = tasks[key]
        if span is not None:
            if "start" not in span or "end" not in span:
                print(
                    f'Error: la tarea "{key}" debe tener "start" y "end" (o ser null).',
                    file=sys.stderr,
                )
                sys.exit(1)
            if span["start"] > span["end"]:
                print(
                    f'Error: en la tarea "{key}", start ({span["start"]}) > end ({span["end"]}).',
                    file=sys.stderr,
                )
                sys.exit(1)
            if span["end"] >= weeks:
                print(
                    f'Error: en la tarea "{key}", end ({span["end"]}) >= weeks ({weeks}).',
                    file=sys.stderr,
                )
                sys.exit(1)

    _require(data, "economic_proposal.development_cost")
    _require(data, "economic_proposal.platform_cost_monthly")
    _require(data, "economic_proposal.support_cost_monthly")


# ── Slides ───────────────────────────────────────────────────────────────────

def build_slide_cover(prs, data):
    slide = prs.slides[0]
    group = slide.shapes[0]
    hero = group.shapes[0]

    title_text = data["cover"]["automation_name"]
    hero.text_frame.clear()
    for i, line in enumerate(title_text.split("\n")):
        p = hero.text_frame.paragraphs[0] if i == 0 else hero.text_frame.add_paragraph()
        r = p.add_run()
        r.text = line
        font_size = 28 if len(line) < 28 else 25
        set_run_font(r, size=font_size, bold=True, color=WHITE)
        p.alignment = PP_ALIGN.LEFT

    slide.shapes[1].text = data["cover"]["year"]
    for p in slide.shapes[1].text_frame.paragraphs:
        for r in p.runs:
            set_run_font(r, size=16, color=WHITE)


def build_slide_understanding(prs, data):
    slide = prs.slides[3]

    objectives = [f"• {x}" for x in data["understanding"]["objectives"]]
    systems = [f"• {x}" for x in data["understanding"]["systems_tools"]]
    set_text_shape(slide.shapes[3], ["Objectives"] + objectives, body_size=15.5)
    set_text_shape(slide.shapes[1], ["Systems/ Tools"] + systems, body_size=15.5)

    # Limpiar imagenes previas antes de insertar la nueva
    for sh in list(slide.shapes):
        if sh.shape_type == 13:  # Picture
            remove_shape(sh)

    flow_path = data["understanding"]["flowchart_image"]
    box = slide.shapes[0]
    pad = Emu(182880)
    img = Image.open(flow_path)
    img_ratio = img.size[0] / img.size[1]
    box_w = box.width - 2 * pad
    box_h = box.height - 2 * pad
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        pic_w = box_w
        pic_h = int(pic_w / img_ratio)
    else:
        pic_h = box_h
        pic_w = int(pic_h * img_ratio)
    pic_left = int(box.left + (box.width - pic_w) / 2)
    pic_top = int(box.top + (box.height - pic_h) / 2)
    slide.shapes.add_picture(flow_path, pic_left, pic_top, width=pic_w, height=pic_h)


def build_slide_complexity(prs, data):
    slide = prs.slides[9]
    table_shape = slide.shapes[1]
    tbl = table_shape.table
    col_widths = [c.width for c in tbl.columns]
    level_map = {"Simple": 1, "Medium": 2, "High": 3, "Very High": 4}
    selected_idx = level_map[data["complexity"]["selected_level"]]

    table_bottom = table_shape.top + table_shape.height
    for sh in list(slide.shapes):
        if (sh.shape_type == 1
                and sh.top >= table_shape.top
                and sh.top + sh.height <= table_bottom):
            remove_shape(sh)

    x = table_shape.left + sum(col_widths[:selected_idx])
    w = col_widths[selected_idx]
    first_row_h = tbl.rows[0].height
    y = table_shape.top + first_row_h
    h = table_shape.height - first_row_h

    highlight = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        x + Emu(15000), y, w - Emu(30000), h,
    )
    highlight.fill.background()
    highlight.line.color.rgb = CYAN
    highlight.line.width = Pt(3.25)


def build_slide_devplan(prs, data):
    slide = prs.slides[10]

    for sh in list(slide.shapes):
        remove_shape(sh)

    title_box = slide.shapes.add_textbox(
        Emu(658368), Emu(283464), Emu(7498080), Emu(566928),
    )
    title_box.text_frame.text = "Development Plan"
    for p in title_box.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            set_run_font(r, size=28, bold=True, color=WHITE)

    tbl_left = 782862
    tbl_top = 1249288
    tbl_width = 10667022
    item_w = 3312060

    weeks = data["development_plan"]["weeks"]
    rows, cols = 17, weeks + 1
    week_w = int((tbl_width - item_w) / weeks)

    tbl_height = 5237018
    tbl_shape = slide.shapes.add_table(
        rows, cols, Emu(tbl_left), Emu(tbl_top), Emu(tbl_width), Emu(tbl_height),
    )
    tbl = tbl_shape.table

    for i, col in enumerate(tbl.columns):
        col.width = Emu(item_w if i == 0 else week_w)

    row_heights = [360218] + [304800] * 16
    for i, row in enumerate(tbl.rows):
        row.height = Emu(row_heights[i])

    tbl_pr = tbl._tbl.tblPr
    style_id_el = tbl_pr.find(qn('a:tblStyleId'))
    if style_id_el is not None:
        tbl_pr.remove(style_id_el)
    for attr in ('bandRow', 'bandCol', 'firstRow', 'lastRow', 'firstCol', 'lastCol'):
        tbl_pr.attrib.pop(attr, None)

    section_rows = {1, 5, 10, 14}
    for r in range(rows):
        for c in range(cols):
            if r == 0:
                fill = NAVY
            elif r in section_rows:
                fill = SECTION_FILL
            else:
                fill = ALT1 if r % 2 == 0 else ALT2
            txt_color = WHITE if r == 0 else DARK_TEXT
            style_cell(
                tbl.cell(r, c), "", fill=fill, font_size=12,
                color=txt_color, bold=(r == 0),
                align=(PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER),
            )

    style_cell(tbl.cell(0, 0), "Item", fill=NAVY, font_size=18,
               bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    for c in range(1, cols):
        style_cell(tbl.cell(0, c), f"W{c - 1}", fill=NAVY, font_size=18,
                   bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    section_titles = {1: "Process Analysis", 5: "Automate", 10: "Test", 14: "Monitor"}
    for r, title in section_titles.items():
        style_cell(tbl.cell(r, 0), "", fill=SECTION_FILL)
        master = tbl.cell(r, 1)
        style_cell(master, title, fill=SECTION_FILL, font_size=16,
                   bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        master.merge(tbl.cell(r, cols - 1))

    labels = {
        2: "Request access and permits to the\nenvironments",
        3: "Process analysis and documentation",
        4: "Process documentation for automation",
        6: "Environment configuration",
        7: "Access validation",
        8: "Automation development",
        9: "Test data generation",
        11: "Unit testing",
        12: "User acceptance test",
        13: "User manual generation",
        15: "Automation rollout",
        16: "Automation monitoring",
    }
    for r, txt in labels.items():
        row_fill = ALT1 if r % 2 == 0 else ALT2
        style_cell(tbl.cell(r, 0), txt, fill=row_fill, font_size=10.5,
                   color=DARK_TEXT, align=PP_ALIGN.LEFT)
        for c in range(1, cols):
            style_cell(tbl.cell(r, c), "", fill=row_fill, font_size=10.5,
                       color=DARK_TEXT, align=PP_ALIGN.CENTER)

    row_to_task = {
        2: "access_permits",
        3: "process_analysis",
        4: "process_doc",
        6: "env_config",
        7: "access_validation",
        8: "automation_dev",
        9: "test_data",
        11: "unit_test",
        12: "uat",
        13: "manual",
        15: "rollout",
        16: "monitoring",
    }

    tasks = data["development_plan"]["tasks"]
    for row_idx, task_key in row_to_task.items():
        span = tasks[task_key]
        if span is None:
            continue
        start_col = 1 + span["start"]
        end_col = 1 + span["end"]
        master = tbl.cell(row_idx, start_col)
        if start_col != end_col:
            master.merge(tbl.cell(row_idx, end_col))
        master.fill.solid()
        master.fill.fore_color.rgb = BLUE


def build_slide_pricing_title(prs):
    slide = prs.slides[11]
    slide.shapes[0].left = Emu(658368)
    slide.shapes[0].top = Emu(310896)
    slide.shapes[0].width = Emu(10424160)
    slide.shapes[0].height = Emu(566928)
    for p in slide.shapes[0].text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            set_run_font(r, size=27, bold=True, color=WHITE)


def build_slide_economic(prs, data):
    slide = prs.slides[13]
    tbl = slide.shapes[1].table

    name = data["cover"]["automation_name"].replace("\n", " ")
    level = data["complexity"]["selected_level"]
    dev_cost = data["economic_proposal"]["development_cost"]
    platform = data["economic_proposal"]["platform_cost_monthly"]
    support = data["economic_proposal"]["support_cost_monthly"]

    tbl.cell(1, 0).text = name
    tbl.cell(1, 1).text = level
    tbl.cell(1, 2).text = currency(dev_cost)
    tbl.cell(1, 3).text = ""
    tbl.cell(2, 3).text = currency(platform)
    tbl.cell(3, 3).text = currency(support)
    tbl.cell(4, 2).text = currency(dev_cost)
    tbl.cell(4, 3).text = currency(platform + support)

    for r in range(5):
        for c in range(4):
            cell = tbl.cell(r, c)
            for p in cell.text_frame.paragraphs:
                if r == 0:
                    p.alignment = PP_ALIGN.CENTER
                elif c == 0:
                    p.alignment = PP_ALIGN.CENTER if r == 4 else PP_ALIGN.LEFT
                else:
                    p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    is_footer = r == 0 or r == 4
                    sz = 15 if (r == 1 and c == 0) else 16
                    clr = WHITE if r == 4 else (DARK_TEXT if r > 0 else WHITE)
                    set_run_font(run, size=sz, bold=is_footer, color=clr)


def build_slides_closing(prs, data):
    year = data["cover"]["year"]
    for idx in [14, 15]:
        slide = prs.slides[idx]
        slide.shapes[1].text = year
        for p in slide.shapes[1].text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                set_run_font(r, size=16, color=WHITE)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Genera una presentacion de estimacion RPA a partir de un JSON de configuracion.",
    )
    parser.add_argument(
        "--config", required=True,
        help="Ruta al archivo JSON de configuracion.",
    )
    parser.add_argument(
        "--out", default="estimacion_output.pptx",
        help="Nombre o ruta del PPTX; relativo -> output/ (default: estimacion_output.pptx).",
    )
    args = parser.parse_args()

    if not os.path.isfile(args.config):
        print(f'Error: no se encuentra el archivo de config "{args.config}".', file=sys.stderr)
        sys.exit(1)

    with open(args.config, encoding="utf-8") as f:
        try:
            data = json.load(f)
        except json.JSONDecodeError as e:
            print(f"Error: el archivo de config no es JSON valido: {e}", file=sys.stderr)
            sys.exit(1)

    validate_config(data)

    template_path = resolve_template_path()
    if not template_path.is_file():
        print(
            f'Error: no se encuentra la plantilla en "{TEMPLATE_DIR}". '
            f"Coloca template.pptx o RPA_Template.pptx en esa carpeta.",
            file=sys.stderr,
        )
        sys.exit(1)

    prs = Presentation(str(template_path))

    build_slide_cover(prs, data)
    build_slide_understanding(prs, data)
    build_slide_complexity(prs, data)
    build_slide_devplan(prs, data)
    build_slide_pricing_title(prs)
    build_slide_economic(prs, data)
    build_slides_closing(prs, data)

    out_path = resolve_output_path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Presentacion generada: {out_path}")


if __name__ == "__main__":
    main()
